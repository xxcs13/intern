import os
import pdfplumber
import pandas as pd
from pptx import Presentation
from dotenv import load_dotenv
from typing import TypedDict, List, Dict, Any

import chromadb
from chromadb.utils import embedding_functions
from openai import OpenAI
from langgraph.graph import StateGraph, END

# ====== Configuration ======
load_dotenv()
openai_api_key = os.getenv("OPENAI_API_KEY")
assert openai_api_key, "Please set OPENAI_API_KEY in .env file"

# ====== State Definition ======
class RAGState(TypedDict):
    pdf_path: str
    pptx_path: str
    xls_path: str
    question: str
    all_chunks: List[Dict[str, Any]]
    collection: Any
    retrieved_chunks: List[Dict[str, Any]]
    answer: str

# ====== Document Parsing Functions ======

def chunk_text(text, chunk_size=1000, overlap=100):
    """Split text into overlapping chunks for better retrieval."""
    chunks = []
    start = 0
    while start < len(text):
        end = start + chunk_size
        chunk = text[start:end]
        chunks.append(chunk)
        start = end - overlap
    return chunks

def parse_pdf(path):
    """Parse PDF file and extract text chunks."""
    chunks = []
    try:
        with pdfplumber.open(path) as pdf:
            for i, page in enumerate(pdf.pages, start=1):
                text = page.extract_text() or ""
                if text.strip():
                    # Split long pages into smaller chunks
                    text_chunks = chunk_text(text.strip())
                    for j, chunk in enumerate(text_chunks):
                        chunks.append({
                            "content": chunk,
                            "metadata": {
                                "source_file": os.path.basename(path), 
                                "type": "pdf", 
                                "page": i,
                                "chunk": j
                            }
                        })
    except Exception as e:
        print(f"Error parsing PDF {path}: {e}")
    return chunks

def parse_pptx(path):
    """Parse PowerPoint file and extract text chunks."""
    chunks = []
    try:
        prs = Presentation(path)
        for i, slide in enumerate(prs.slides, start=1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text.append(shape.text)
                if shape.has_table:
                    for row in shape.table.rows:
                        slide_text.append(" | ".join([cell.text for cell in row.cells]))
            text = "\n".join(slide_text)
            if text.strip():
                chunks.append({
                    "content": text.strip(),
                    "metadata": {
                        "source_file": os.path.basename(path), 
                        "type": "pptx", 
                        "slide": i
                    }
                })
    except Exception as e:
        print(f"Error parsing PPTX {path}: {e}")
    return chunks

def parse_xls(path):
    """Parse Excel file and extract data chunks."""
    chunks = []
    try:
        xls = pd.ExcelFile(path)
        for sheet in xls.sheet_names:
            df = xls.parse(sheet).astype(str)
            # Combine every N rows into one chunk
            chunk_size = 20
            for i in range(0, len(df), chunk_size):
                chunk_df = df.iloc[i:i+chunk_size]
                text = "\n".join([", ".join(row) for row in chunk_df.values])
                if text.strip():
                    chunks.append({
                        "content": text.strip(),
                        "metadata": {
                            "source_file": os.path.basename(path), 
                            "type": "xls", 
                            "sheet": sheet, 
                            "rows": f"{i+1}-{i+len(chunk_df)}"
                        }
                    })
    except Exception as e:
        print(f"Error parsing XLS {path}: {e}")
    return chunks

# ====== Embedding & Vector Database ======
def build_or_load_vector_db(all_chunks, db_path="chroma_db"):
    """Build ChromaDB vector database from document chunks or load existing one."""
    try:
        ef = embedding_functions.OpenAIEmbeddingFunction(
            api_key=openai_api_key,
            model_name="text-embedding-3-small"
        )
        client = chromadb.PersistentClient(path=db_path)
        
        # Check if collection already exists
        try:
            collection = client.get_collection("rag_chunks", embedding_function=ef)
            print(f"Found existing vector database with {collection.count()} chunks")
            return collection
        except Exception:
            # Collection doesn't exist, create it
            print("No existing vector database found, creating new one...")
            pass
        
        # Create new collection
        collection = client.create_collection("rag_chunks", embedding_function=ef)
        
        # Add chunks to collection
        for idx, chunk in enumerate(all_chunks):
            content = chunk["content"]
            meta = chunk["metadata"]
            meta = dict(meta)
            meta["chunk_id"] = idx
            
            # Create unique ID for each chunk
            source_id = meta.get('page', meta.get('slide', meta.get('sheet', '')))
            chunk_id = f"{meta['source_file']}_{meta['type']}_{source_id}_{idx}"
            
            collection.add(
                documents=[content],
                metadatas=[meta],
                ids=[chunk_id]
            )
        
        print(f"Successfully created and added {len(all_chunks)} chunks to vector database")
        return collection
    except Exception as e:
        print(f"Error with vector database: {e}")
        raise

def force_rebuild_vector_db(all_chunks, db_path="chroma_db"):
    """Force rebuild the vector database (useful for updates)."""
    try:
        ef = embedding_functions.OpenAIEmbeddingFunction(
            api_key=openai_api_key,
            model_name="text-embedding-3-small"
        )
        client = chromadb.PersistentClient(path=db_path)
        
        # Clean up existing collection
        try:
            client.delete_collection("rag_chunks")
            print("Deleted existing vector database")
        except Exception:
            pass
        
        # Create new collection
        collection = client.create_collection("rag_chunks", embedding_function=ef)
        
        # Add chunks to collection
        for idx, chunk in enumerate(all_chunks):
            content = chunk["content"]
            meta = chunk["metadata"]
            meta = dict(meta)
            meta["chunk_id"] = idx
            
            # Create unique ID for each chunk
            source_id = meta.get('page', meta.get('slide', meta.get('sheet', '')))
            chunk_id = f"{meta['source_file']}_{meta['type']}_{source_id}_{idx}"
            
            collection.add(
                documents=[content],
                metadatas=[meta],
                ids=[chunk_id]
            )
        
        print(f"Successfully rebuilt vector database with {len(all_chunks)} chunks")
        return collection
    except Exception as e:
        print(f"Error rebuilding vector database: {e}")
        raise

# ====== RAG Query & LLM Generation ======
def query_rag(question, collection, top_k=5):
    """Query the vector database for relevant chunks."""
    try:
        results = collection.query(
            query_texts=[question],
            n_results=top_k
        )
        retrieved_chunks = []
        for doc, meta in zip(results["documents"][0], results["metadatas"][0]):
            retrieved_chunks.append({
                "content": doc,
                "metadata": meta
            })
        return retrieved_chunks
    except Exception as e:
        print(f"Error querying vector database: {e}")
        return []

def llm_generate_answer(question, context_chunks):
    """Generate answer using OpenAI LLM based on retrieved context."""
    try:
        client = OpenAI(api_key=openai_api_key)
        context = "\n\n".join([
            f"[Source: {chunk['metadata']['source_file']} - {chunk['metadata']['type']} - Page/Slide: {chunk['metadata'].get('page', chunk['metadata'].get('slide', chunk['metadata'].get('sheet', 'N/A')))}]\n{chunk['content']}" 
            for chunk in context_chunks
        ])
        
        prompt = f"""You are a knowledgeable assistant. Please answer the question based on the following context:

Context:
{context}

Question: {question}

Please provide a direct answer and mention the source files, pages/slides/sheets you used for the answer."""

        response = client.chat.completions.create(
            model="gpt-3.5-turbo",
            messages=[{"role": "user", "content": prompt}],
            temperature=0.1,
            max_tokens=800
        )
        return response.choices[0].message.content
    except Exception as e:
        print(f"Error generating answer: {e}")
        return f"Error generating answer: {e}"

# ====== LangGraph Workflow Nodes ======

def parse_all_node(state: RAGState) -> RAGState:
    """Parse all documents (PDF, PPTX, XLS) and extract chunks."""
    print("Parsing documents...")
    pdf_chunks = parse_pdf(state['pdf_path'])
    pptx_chunks = parse_pptx(state['pptx_path'])
    xls_chunks = parse_xls(state['xls_path'])
    all_chunks = pdf_chunks + pptx_chunks + xls_chunks
    print(f"Total chunks extracted: {len(all_chunks)}")
    return {"all_chunks": all_chunks}

def embedding_node(state: RAGState) -> RAGState:
    """Build vector database with embeddings or load existing one."""
    print("Loading or building vector database...")
    collection = build_or_load_vector_db(state["all_chunks"])
    return {"collection": collection}

def query_node(state: RAGState) -> RAGState:
    """Query vector database for relevant chunks."""
    print("Querying vector database...")
    question = state["question"]
    collection = state["collection"]
    retrieved_chunks = query_rag(question, collection)
    print(f"Retrieved {len(retrieved_chunks)} relevant chunks")

    updated_state = state.copy()
    updated_state["retrieved_chunks"] = retrieved_chunks
    return updated_state

def rag_node(state: RAGState) -> RAGState:
    """Generate answer using LLM based on retrieved chunks."""
    print("Generating answer...")
    question = state["question"]
    chunks = state["retrieved_chunks"]
    answer = llm_generate_answer(question, chunks)

    updated_state = state.copy()
    updated_state["answer"] = answer
    return updated_state

def log_node(state: RAGState) -> RAGState:
    """Log Q&A session to CSV file."""
    print("Logging Q&A session...")
    row = {
        "question": state["question"],
        "answer": state["answer"],
        "retrieved_chunks": str(state["retrieved_chunks"])
    }
    df = pd.DataFrame([row])
    if not os.path.exists("rag_qa_log.csv"):
        df.to_csv("rag_qa_log.csv", index=False)
    else:
        df.to_csv("rag_qa_log.csv", mode="a", header=False, index=False)
    return state

# ====== Build LangGraph Workflow ======
workflow = StateGraph(RAGState)
workflow.add_node("parse_all", parse_all_node)
workflow.add_node("embedding", embedding_node)
workflow.add_node("query", query_node)
workflow.add_node("rag", rag_node)
workflow.add_node("log", log_node)

# Define workflow connections
workflow.set_entry_point("parse_all")
workflow.add_edge("parse_all", "embedding")
workflow.add_edge("embedding", "query")
workflow.add_edge("query", "rag")
workflow.add_edge("rag", "log")
workflow.add_edge("log", END)

# Compile the workflow
graph = workflow.compile()

# ====== Main Execution ======
if __name__ == "__main__":
    print("Starting RAG Q&A System...")
    print("=" * 50)
    
    # Define the initial state for document processing
    base_state = {
        "pdf_path": "./tsmc_2024_yearly report.pdf",
        "pptx_path": "./無備忘錄版_digitimes_2025年全球AI伺服器出貨將達181萬台　高階機種採購不再集中於四大CSP.pptx",
        "xls_path": "./excel_FS-Consolidated_1Q25.xls",
        "question": "",  # Will be set for each question
        "all_chunks": [],
        "collection": None,
        "retrieved_chunks": [],
        "answer": ""
    }
    
    # Check if user wants to rebuild database
    rebuild_db = input("是否要重新建構向量數據庫？(y/N): ").strip().lower()
    
    # First, process all documents once (parse and build vector database)
    print("Initializing system - processing documents...")
    try:
        # Parse all documents and build vector database
        temp_state = base_state.copy()
        temp_state["question"] = "initialization"
        
        # Check if we need to rebuild or can use existing database
        if rebuild_db == 'y':
            print("Forcing rebuild of vector database...")
            temp_state = parse_all_node(temp_state)
            collection = force_rebuild_vector_db(temp_state["all_chunks"])
            all_chunks = temp_state["all_chunks"]
        else:
            # Try to load existing database first
            try:
                ef = embedding_functions.OpenAIEmbeddingFunction(
                    api_key=openai_api_key,
                    model_name="text-embedding-3-small"
                )
                client = chromadb.PersistentClient(path="chroma_db")
                collection = client.get_collection("rag_chunks", embedding_function=ef)
                print(f"Loaded existing vector database with {collection.count()} chunks")
                # We still need to parse documents to get all_chunks for potential use
                temp_state = parse_all_node(temp_state)
                all_chunks = temp_state["all_chunks"]
            except Exception:
                print("No existing database found, creating new one...")
                temp_state = parse_all_node(temp_state)
                temp_state = embedding_node(temp_state)
                collection = temp_state["collection"]
                all_chunks = temp_state["all_chunks"]
        
        print("System initialized successfully!")
        print("=" * 50)
        
        # Interactive Q&A loop
        question_count = 0
        while True:
            print(f"\n問題 {question_count + 1}:")
            question = input("請輸入您的問題 (輸入 'quit' 或 'exit' 結束): ").strip()
            
            if question.lower() in ['quit', 'exit', 'q']:
                print("感謝使用！再見！")
                break
            
            if not question:
                print("請輸入有效的問題。")
                continue
            
            try:
                # Create state for this question
                current_state = {
                    "pdf_path": base_state["pdf_path"],
                    "pptx_path": base_state["pptx_path"],
                    "xls_path": base_state["xls_path"],
                    "question": question,
                    "all_chunks": all_chunks,
                    "collection": collection,
                    "retrieved_chunks": [],
                    "answer": ""
                }
                
                # Execute query and answer generation

                current_state = query_node(current_state)
                current_state = rag_node(current_state)
                
                # Log the Q&A
                log_node(current_state)
                
                # Display results
                print("\n" + "=" * 20 + " 問題 " + "=" * 20)
                print(question)
                print("\n" + "=" * 20 + " 回答 " + "=" * 20)
                print(current_state["answer"])
                print("\n" + "=" * 50)
                
                question_count += 1
                
            except Exception as e:
                print(f"處理問題時發生錯誤: {e}")
                import traceback
                traceback.print_exc()
                print("請嘗試重新輸入問題。")
        
    except Exception as e:
        print(f"初始化系統時發生錯誤: {e}")
        import traceback
        traceback.print_exc()
