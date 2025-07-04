#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Extract all tables and charts from a PowerPoint file and save each as a CSV.

Requirements:
    pip install python-pptx pandas
"""

import os
from pathlib import Path
import pandas as pd
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def extract_table_shape(table, csv_path: Path) -> None:
    """Save a PowerPoint Table shape to CSV."""
    data = [
        [cell.text_frame.text.replace("\n", " ").strip() for cell in row.cells]
        for row in table.rows
    ]
    pd.DataFrame(data).to_csv(csv_path, index=False, header=False, encoding="utf-8-sig")


def extract_chart_shape(chart, csv_path: Path) -> None:
    """Save chart data (categories & series values) to CSV."""
    # Attempt to read categories; fallback to index if absent
    categories = []
    try:
        categories = [c.label if hasattr(c, "label") else str(c) for c in chart.plots[0].categories]  # type: ignore
    except Exception:
        pass

    rows, header = [], []
    header = ["" if categories else "Index"] + [s.name for s in chart.series]

    # Determine row count (length of longest series)
    row_count = max(len(s.values) for s in chart.series)

    for i in range(row_count):
        cat = categories[i] if categories else i
        row = [cat]
        for series in chart.series:
            value = series.values[i]
            # Each value can be a numeric element or string; handle both
            row.append(value.value if hasattr(value, "value") else value)
        rows.append(row)

    # Create DataFrame and transpose it (swap rows and columns)
    df = pd.DataFrame(rows, columns=header)
    df_transposed = df.T
    df_transposed.to_csv(csv_path, header=False, encoding="utf-8-sig")


def extract_ppt_tables_to_csv(
    ppt_path: str,
    output_dir: str | None = None,
    slide_range: tuple[int, int] | None = None,
) -> list[Path]:
    """
    Extract all Table shapes and Chart objects in a PPTX to CSV.

    Args:
        ppt_path: Path to the .pptx file.
        output_dir: Directory to save CSV files; defaults to PPT's directory.
        slide_range: (start, end) 1-based inclusive slide numbers to process. None → all.

    Returns:
        List of generated CSV Path objects.
    """
    prs = Presentation(ppt_path)
    out_dir = Path(output_dir) if output_dir else Path(ppt_path).parent / "ppt"
    out_dir.mkdir(parents=True, exist_ok=True)

    csv_paths: list[Path] = []
    tbl_idx = chart_idx = 0

    for slide_num, slide in enumerate(prs.slides, start=1):
        # Skip slides outside desired range
        if slide_range and not (slide_range[0] <= slide_num <= slide_range[1]):
            continue

        for shape in slide.shapes:
            # ---- Table shape ----
            if shape.has_table:
                tbl_idx += 1
                csv_file = out_dir / f"table_{tbl_idx}_slide{slide_num}.csv"
                extract_table_shape(shape.table, csv_file)
                csv_paths.append(csv_file)

            # ---- Chart shape ----
            elif shape.has_chart:
                chart_idx += 1
                csv_file = out_dir / f"chart_{chart_idx}_slide{slide_num}.csv"
                extract_chart_shape(shape.chart, csv_file)
                csv_paths.append(csv_file)


    return csv_paths


if __name__ == "__main__":
    # --------- User settings ---------
    ppt_path = "無備忘錄版_digitimes_2025年全球AI伺服器出貨將達181萬台　高階機種採購不再集中於四大CSP.pptx"
    output_directory = None        # e.g. "./csv_output"
    slide_range = None             # e.g. (4, 10)  → only slides 4-10
    # ---------------------------------

    generated_files = extract_ppt_tables_to_csv(
        ppt_path=ppt_path,
        output_dir=output_directory,
        slide_range=slide_range,
    )

    print("CSV files generated:")
    for p in generated_files:
        print(" •", p)
