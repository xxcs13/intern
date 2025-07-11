import os
import pdfplumber
import pandas as pd
from pptx import Presentation
from dotenv import load_dotenv
from typing import List, Dict, Any

from langchain_community.vectorstores import Chroma
from langchain_openai import OpenAIEmbeddings, ChatOpenAI
from langchain.schema.document import Document

from langgraph.graph import StateGraph, END


from dataclasses import dataclass, field

@dataclass
class GraphState:
    docs: List[Document] = field(default_factory=list)
    vectorstore: Any = None
    question: str = ""
    retrieved_docs: List[Document] = field(default_factory=list)
    answer: str = ""


load_dotenv()
openai_api_key = os.getenv("OPENAI_API_KEY")
assert openai_api_key, "請在 .env 設定 OPENAI_API_KEY"


def parse_pdf(path: str) -> List[Document]:
    docs = []
    with pdfplumber.open(path) as pdf:
        for i, page in enumerate(pdf.pages, start=1):
            text = page.extract_text() or ""
            if text.strip():
                docs.append(Document(
                    page_content=text.strip(),
                    metadata={"source_file": os.path.basename(path), "type": "pdf", "page": i}
                ))
    return docs

def parse_pptx(path: str) -> List[Document]:
    docs = []
    prs = Presentation(path)
    for i, slide in enumerate(prs.slides, start=1):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                slide_text.append(shape.text)
            if shape.has_table:
                for row in shape.table.rows:
                    slide_text.append(" | ".join([cell.text for cell in row.cells]))
        text = "\n".join(slide_text)
        if text.strip():
            docs.append(Document(
                page_content=text.strip(),
                metadata={"source_file": os.path.basename(path), "type": "pptx", "slide": i}
            ))
    return docs

def parse_xls(path: str) -> List[Document]:
    docs = []
    xls = pd.ExcelFile(path)
    for sheet in xls.sheet_names:
        df = xls.parse(sheet).astype(str)
        chunk_size = 20
        for i in range(0, len(df), chunk_size):
            chunk_df = df.iloc[i:i+chunk_size]
            text = "\n".join([", ".join(row) for row in chunk_df.values])
            if text.strip():
                docs.append(Document(
                    page_content=text.strip(),
                    metadata={"source_file": os.path.basename(path), "type": "xls", "sheet": sheet, "rows": f"{i+1}-{i+len(chunk_df)}"}
                ))
    return docs


def ingest_node(state: GraphState) -> GraphState:

    pdf_docs = parse_pdf(state.docs[0])   # state.docs[0] 是 pdf_path
    pptx_docs = parse_pptx(state.docs[1]) # state.docs[1] 是 pptx_path
    xls_docs = parse_xls(state.docs[2])   # state.docs[2] 是 xls_path
    docs = pdf_docs + pptx_docs + xls_docs
    return GraphState(
        docs=docs,
        vectorstore=None,
        question=state.question
    )

def embed_node(state: GraphState) -> GraphState:
    embeddings = OpenAIEmbeddings(openai_api_key=openai_api_key)
    vectorstore = Chroma.from_documents(
        state.docs,
        embedding=embeddings,
        persist_directory="chroma_db_rag"
    )
    return GraphState(
        docs=state.docs,
        vectorstore=vectorstore,
        question=state.question
    )

def retrieval_node(state: GraphState) -> GraphState:
    docs = state.vectorstore.similarity_search(state.question, k=4)
    return GraphState(
        docs=state.docs,
        vectorstore=state.vectorstore,
        question=state.question,
        retrieved_docs=docs
    )

def rag_node(state: GraphState) -> GraphState:
    llm = ChatOpenAI(model="gpt-3.5-turbo", temperature=0, openai_api_key=openai_api_key)
    context = "\n\n".join([
        f"【來源:{doc.metadata}】\n{doc.page_content}" for doc in state.retrieved_docs
    ])
    prompt = (
        f"You are a knowledgeable assistant. Please answer the question based on the following context:\n\n"
        f"{context}\n\n"
        f"問題：{state.question}\n\n"
        f"請直接作答，並標明所用的來源檔案和頁碼/slide/sheet。"
    )
    answer = llm.invoke(prompt)
    return GraphState(
        docs=state.docs,
        vectorstore=state.vectorstore,
        question=state.question,
        retrieved_docs=state.retrieved_docs,
        answer=answer.content if hasattr(answer, "content") else str(answer)
    )

def log_node(state: GraphState) -> GraphState:
    row = {
        "question": state.question,
        "answer": state.answer,
        "retrieved_docs": str([doc.metadata for doc in state.retrieved_docs])
    }
    df = pd.DataFrame([row])
    if not os.path.exists("rag_qa_log.csv"):
        df.to_csv("rag_qa_log.csv", index=False)
    else:
        df.to_csv("rag_qa_log.csv", mode="a", header=False, index=False)
    return state

workflow = StateGraph(GraphState)
workflow.add_node("ingest", ingest_node)
workflow.add_node("embed", embed_node)
workflow.add_node("retrieval", retrieval_node)
workflow.add_node("rag", rag_node)
workflow.add_node("log", log_node)

workflow.set_entry("ingest")
workflow.connect("ingest", "embed")
workflow.connect("embed", "retrieval")
workflow.connect("retrieval", "rag")
workflow.connect("rag", "log")
workflow.connect("log", END)
graph = workflow.compile()


if __name__ == "__main__":

    PDF_PATH = "./tsmc_2024_yearly report.pdf"
    PPTX_PATH = "./無備忘錄版_digitimes_2025年全球AI伺服器出貨將達181萬台　高階機種採購不再集中於四大CSP.pptx"
    XLS_PATH = "./excel_FS-Consolidated_1Q25.xls"
    QUESTION = "2025年全球AI伺服器的市場規模是多少？"

    initial_state = GraphState(
        docs=[PDF_PATH, PPTX_PATH, XLS_PATH],
        question=QUESTION
    )
    result = graph.invoke(initial_state)
    print("==== 問題 ====")
    print(QUESTION)
    print("==== 答案 ====")
    print(result.answer)
